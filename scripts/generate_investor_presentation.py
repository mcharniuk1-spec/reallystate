#!/usr/bin/env python3
"""Generate YC-style investor presentation artifacts (PDF + PPTX).

Outputs:
- output/pdf/investor-presentation-<date>.pdf
- docs/exports/investor-presentation-<date>.pptx
"""

from __future__ import annotations

import sys
from dataclasses import dataclass
from datetime import date
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
TMP = ROOT / "tmp" / "presentations" / "charts"
OUT_PDF_DIR = ROOT / "output" / "pdf"
OUT_PPTX_DIR = ROOT / "docs" / "exports"
OUT_NOTES_DIR = ROOT / "docs" / "exports"


def _ensure_pdf_deps() -> None:
    missing: list[str] = []
    try:
        import matplotlib  # noqa: F401
    except ImportError:
        missing.append("matplotlib")
    try:
        import reportlab  # noqa: F401
    except ImportError:
        missing.append("reportlab")
    if missing:
        print(f"[ERROR] Missing required package(s): {', '.join(missing)}")
        print(f"Install with: python3 -m pip install {' '.join(missing)}")
        sys.exit(1)


def _pptx_available() -> bool:
    try:
        import pptx  # noqa: F401
        return True
    except ImportError:
        return False


_ensure_pdf_deps()

import matplotlib  # noqa: E402

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from reportlab.lib import colors  # noqa: E402
from reportlab.lib.pagesizes import A4, landscape  # noqa: E402
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet  # noqa: E402
from reportlab.lib.units import cm  # noqa: E402
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer  # noqa: E402

if _pptx_available():
    from pptx import Presentation  # noqa: E402
    from pptx.enum.text import PP_ALIGN  # noqa: E402
    from pptx.util import Inches, Pt  # noqa: E402


BRAND_BLUE = "#2457d6"
BRAND_GREEN = "#17a672"
BRAND_AMBER = "#e09f3e"
BRAND_RED = "#d64545"
BRAND_DARK = "#0f1b2d"

W, H = landscape(A4)
MARGIN = 1.5 * cm


@dataclass
class Slide:
    title: str
    paragraph: str
    chart_key: str | None
    sources: str


SLIDES: list[Slide] = [
    Slide(
        "Bulgaria Real Estate Platform",
        (
            "We are building the first unified property marketplace for Bulgaria that combines multi-source listing aggregation, "
            "deduplication, map-first search, and an AI property assistant. The immediate wedge is Varna and the coastal corridor, "
            "where price growth, construction intensity, and short-term rental demand create a high-velocity market with clear monetization."
        ),
        "title_kpi",
        "Sources: NSI, Eurostat, project source registry (2024-2026).",
    ),
    Slide(
        "Problem: Fragmented Supply, Broken Discovery",
        (
            "Buyers and renters must manually check many disconnected portals, agencies, and social channels to see a representative set of supply. "
            "This creates search fatigue, duplicate listings, and delayed decisions while agencies and owners lose qualified demand because discovery is fragmented."
        ),
        "source_mix",
        "Sources: Project source matrix, Bulgarian portals audit (2026-04).",
    ),
    Slide(
        "Why Now: Market Momentum Is Exceptional",
        (
            "Bulgaria's residential market is expanding faster than the EU average, with coastal cities showing sustained momentum and strong transaction activity. "
            "This timing allows a modern product to establish distribution before incumbents close UX and data-quality gaps."
        ),
        "market_growth",
        "Sources: Eurostat HPI, national market reviews, regional agency reports.",
    ),
    Slide(
        "Beachhead Market: Varna + Burgas",
        (
            "Varna is the MVP focus because it combines high listing churn, strong buyer intent, and compelling map-based discovery use cases. "
            "Burgas is a natural second coastal expansion to scale supply depth and improve model quality for price and demand prediction."
        ),
        "varna_burgas",
        "Sources: Regional price and construction snapshots (2024-2026).",
    ),
    Slide(
        "Market Size: TAM, SAM, SOM",
        (
            "The addressable revenue pool supports a venture-scale business if execution compounds from coastal leadership to national coverage. "
            "Our SOM assumptions are grounded in phased rollout, controlled CAC, and progressive source onboarding with legal/compliance gates."
        ),
        "tam_sam_som",
        "Sources: Internal unit economics model and market sizing assumptions.",
    ),
    Slide(
        "Business Model and Unit Economics",
        (
            "Revenue is diversified across agency subscriptions, premium listings, lead products, analytics, and future AI add-ons. "
            "The model targets strong gross margins by keeping ingestion and ranking systems highly automated while preserving compliance controls."
        ),
        "unit_econ",
        "Sources: Internal pricing model, planned SKU structure, cost baseline.",
    ),
    Slide(
        "Go-To-Market and Expansion",
        (
            "The launch motion starts with focused inventory quality and buyer-side utility in Varna, then expands to Burgas and additional cities. "
            "Each stage increases defensibility through better coverage, better ranking, and higher switching costs for agencies and power users."
        ),
        "gtm",
        "Sources: Execution roadmap in PLAN.md and TASKS.md.",
    ),
    Slide(
        "Product Experience: Search, Map, AI Chat",
        (
            "The core experience is map-first with rich filters, canonical property pages, and an AI chat layer that helps users narrow options quickly. "
            "This creates a faster decision loop than static listing portals and supports future conversational commerce flows."
        ),
        "product_status",
        "Sources: Product UX structure and architecture execution guide.",
    ),
    Slide(
        "Execution Progress and Delivery Velocity",
        (
            "The program is organized into parallel specialist lanes with explicit dependencies and verification gates. "
            "This structure keeps engineering throughput high while preserving quality on data, compliance, and product experience."
        ),
        "execution_progress",
        "Sources: docs/agents/TASKS.md, dashboard exports, journey logs.",
    ),
    Slide(
        "Defensibility and Moat",
        (
            "Our moat compounds from canonical data quality, cross-source deduplication, legal-safe ingestion workflows, and map/chat UX differentiation. "
            "As data density and user interaction increase, search relevance and lead quality improve, reinforcing the business flywheel."
        ),
        "moat",
        "Sources: Platform architecture and competitive benchmarking assumptions.",
    ),
    Slide(
        "Risks and Mitigations",
        (
            "Primary risks include source volatility, legal constraints, and execution complexity. "
            "Mitigation is built into the operating model through source-tier controls, fixture-first testing, auditability, and staged geographic rollouts."
        ),
        "risk",
        "Sources: Source registry legal/risk/access gates and QA protocol.",
    ),
    Slide(
        "The Ask",
        (
            "We are preparing a public startup-grade launch and fundraising story with a clear sequence: finish stage-one scraping quality gates, "
            "ship Varna map-plus-chat MVP, and scale monetization by expanding verified supply and lead conversion outcomes."
        ),
        "ask",
        "Sources: Current roadmap and investor model assumptions.",
    ),
]


def _save_chart(fig: plt.Figure, name: str) -> Path:
    TMP.mkdir(parents=True, exist_ok=True)
    out = TMP / f"{name}.png"
    fig.savefig(out, dpi=170, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out


def _chart_title_kpi() -> Path:
    fig, ax = plt.subplots(figsize=(10, 4))
    ax.axis("off")
    metrics = [
        ("30+", "Active + planned sources"),
        ("6", "Parallel specialist lanes"),
        ("Varna", "MVP geography"),
        ("YC", "Startup-grade pitch track"),
    ]
    x = [0.12, 0.37, 0.62, 0.87]
    for idx, (value, label) in enumerate(metrics):
        ax.text(x[idx], 0.68, value, ha="center", va="center", fontsize=30, color=BRAND_BLUE, fontweight="bold")
        ax.text(x[idx], 0.42, label, ha="center", va="center", fontsize=11, color=BRAND_DARK)
    return _save_chart(fig, "title_kpi")


def _chart_source_mix() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    labels = ["Portals", "Classifieds", "Agencies", "International", "Vendor/API", "Social"]
    values = [30, 22, 25, 8, 10, 5]
    colors_m = [BRAND_BLUE, "#6f42c1", BRAND_GREEN, "#0aa2c0", BRAND_AMBER, "#8f9aa8"]
    ax.pie(values, labels=labels, autopct="%1.0f%%", colors=colors_m, startangle=100, textprops={"fontsize": 10})
    ax.set_title("Listing distribution by source class", fontsize=13, fontweight="bold")
    return _save_chart(fig, "source_mix")


def _chart_market_growth() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    cities = ["Sofia", "Varna", "Burgas", "Plovdiv", "EU avg"]
    yoy = [18, 17, 19, 13, 3]
    bars = ax.barh(cities, yoy, color=[BRAND_BLUE, BRAND_GREEN, BRAND_AMBER, "#7f8c8d", "#c4c9cf"])
    for b, v in zip(bars, yoy):
        ax.text(v + 0.3, b.get_y() + b.get_height() / 2, f"{v}%", va="center", fontsize=10)
    ax.set_xlim(0, 24)
    ax.set_xlabel("YoY growth (%)")
    ax.set_title("Residential momentum: Bulgaria cities vs EU baseline", fontsize=13, fontweight="bold")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart(fig, "market_growth")


def _chart_varna_burgas() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    cats = ["Avg price EUR/m2", "YoY growth %", "ADR EUR/night", "Occupancy %"]
    varna = [1900, 17, 65, 59]
    burgas = [1700, 19, 58, 50]
    x = range(len(cats))
    ax.bar([i - 0.16 for i in x], varna, width=0.32, color=BRAND_BLUE, label="Varna")
    ax.bar([i + 0.16 for i in x], burgas, width=0.32, color=BRAND_GREEN, label="Burgas")
    ax.set_xticks(list(x), cats, fontsize=9)
    ax.set_title("Coastal benchmark metrics", fontsize=13, fontweight="bold")
    ax.legend()
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart(fig, "varna_burgas")


def _chart_tam_sam_som() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    segments = ["TAM", "SAM", "SOM (Y3)"]
    vals = [74.7, 22.6, 5.6]
    cols = [BRAND_BLUE, BRAND_GREEN, BRAND_AMBER]
    bars = ax.bar(segments, vals, color=cols)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + 1.0, f"EUR {v:.1f}M", ha="center", fontsize=10, fontweight="bold")
    ax.set_ylim(0, 84)
    ax.set_ylabel("Annual revenue pool (EUR millions)")
    ax.set_title("Market sizing framework", fontsize=13, fontweight="bold")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart(fig, "tam_sam_som")


def _chart_unit_econ() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    revenue = [0.7, 2.3, 4.5]
    cost = [0.14, 0.4, 0.8]
    years = ["Y1", "Y2", "Y3"]
    x = range(len(years))
    ax.bar([i - 0.17 for i in x], revenue, width=0.34, label="Revenue", color=BRAND_BLUE)
    ax.bar([i + 0.17 for i in x], cost, width=0.34, label="Operating cost", color="#b6c0ce")
    ax.set_xticks(list(x), years)
    ax.set_ylabel("EUR millions")
    ax.set_title("Revenue vs operating cost (conservative path)", fontsize=13, fontweight="bold")
    ax.legend()
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart(fig, "unit_econ")


def _chart_gtm() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    phases = ["MVP", "Coastal", "National", "Scale"]
    targets = [5, 15, 50, 100]
    ax.plot(phases, targets, marker="o", linewidth=3, color=BRAND_BLUE)
    for p, t in zip(phases, targets):
        ax.text(p, t + 2, f"{t}K MAU", ha="center", fontsize=10)
    ax.set_ylim(0, 115)
    ax.set_ylabel("Target monthly active users (K)")
    ax.set_title("Go-to-market growth milestones", fontsize=13, fontweight="bold")
    ax.grid(axis="y", linestyle="--", alpha=0.35)
    return _save_chart(fig, "gtm")


def _chart_product_status() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    features = ["Ingestion", "Canonical DB", "Map UX", "AI Chat", "CRM", "Publishing"]
    progress = [72, 65, 52, 38, 44, 20]
    colors_f = [BRAND_GREEN if p >= 65 else BRAND_AMBER if p >= 35 else BRAND_RED for p in progress]
    bars = ax.barh(features, progress, color=colors_f)
    for b, p in zip(bars, progress):
        ax.text(p + 1, b.get_y() + b.get_height() / 2, f"{p}%", va="center", fontsize=10)
    ax.set_xlim(0, 100)
    ax.set_xlabel("Delivery progress (%)")
    ax.set_title("Product capability readiness", fontsize=13, fontweight="bold")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    return _save_chart(fig, "product_status")


def _chart_execution_progress() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    lanes = ["Backend", "Scraper T1/T2", "Scraper T3", "Social", "UX/UI", "Debugger"]
    pct = [42, 48, 35, 30, 40, 28]
    ax.bar(lanes, pct, color=[BRAND_BLUE, BRAND_GREEN, BRAND_AMBER, "#6f42c1", "#0aa2c0", "#7f8c8d"])
    ax.set_ylim(0, 100)
    ax.set_ylabel("Progress (%)")
    ax.set_title("Parallel lane execution status", fontsize=13, fontweight="bold")
    ax.tick_params(axis="x", rotation=18)
    return _save_chart(fig, "execution_progress")


def _chart_moat() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    factors = ["Coverage", "Dedup quality", "Compliance", "Map UX", "AI search"]
    values = [8.5, 8.0, 8.8, 7.9, 7.5]
    ax.plot(factors, values, marker="o", linewidth=3, color=BRAND_BLUE)
    ax.fill_between(factors, values, [0] * len(values), color="#2457d61f")
    ax.set_ylim(0, 10)
    ax.set_ylabel("Relative strength (0-10)")
    ax.set_title("Defensibility profile", fontsize=13, fontweight="bold")
    ax.grid(axis="y", linestyle="--", alpha=0.3)
    return _save_chart(fig, "moat")


def _chart_risk() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    risks = ["Source blocking", "Legal change", "Execution delay", "Data quality", "Adoption lag"]
    severity = [6, 5, 6, 4, 5]
    mitigation = [8, 7, 7, 8, 6]
    x = range(len(risks))
    ax.bar([i - 0.16 for i in x], severity, width=0.32, color=BRAND_RED, label="Risk")
    ax.bar([i + 0.16 for i in x], mitigation, width=0.32, color=BRAND_GREEN, label="Mitigation strength")
    ax.set_xticks(list(x), risks, rotation=20)
    ax.set_ylim(0, 10)
    ax.set_title("Risk and mitigation balance", fontsize=13, fontweight="bold")
    ax.legend()
    return _save_chart(fig, "risk")


def _chart_ask() -> Path:
    fig, ax = plt.subplots(figsize=(9, 4.5))
    ax.axis("off")
    bullets = [
        "1) Finish stage-one scraping quality gate",
        "2) Launch Varna map + AI chat MVP",
        "3) Expand verified supply and monetization",
        "4) Scale to additional Bulgarian cities",
    ]
    for i, txt in enumerate(bullets):
        ax.text(0.05, 0.82 - i * 0.2, txt, fontsize=14, color=BRAND_DARK)
    ax.text(0.05, 0.08, "Public startup presentation track aligned to YC expectations.", fontsize=12, color=BRAND_BLUE)
    return _save_chart(fig, "ask")


CHART_BUILDERS = {
    "title_kpi": _chart_title_kpi,
    "source_mix": _chart_source_mix,
    "market_growth": _chart_market_growth,
    "varna_burgas": _chart_varna_burgas,
    "tam_sam_som": _chart_tam_sam_som,
    "unit_econ": _chart_unit_econ,
    "gtm": _chart_gtm,
    "product_status": _chart_product_status,
    "execution_progress": _chart_execution_progress,
    "moat": _chart_moat,
    "risk": _chart_risk,
    "ask": _chart_ask,
}


def _build_chart_files() -> dict[str, Path]:
    out: dict[str, Path] = {}
    for key, builder in CHART_BUILDERS.items():
        out[key] = builder()
    return out


def build_pdf(chart_paths: dict[str, Path]) -> Path:
    OUT_PDF_DIR.mkdir(parents=True, exist_ok=True)
    out_pdf = OUT_PDF_DIR / f"investor-presentation-{date.today().isoformat()}.pdf"
    doc = SimpleDocTemplate(
        str(out_pdf),
        pagesize=landscape(A4),
        leftMargin=MARGIN,
        rightMargin=MARGIN,
        topMargin=1.2 * cm,
        bottomMargin=1.0 * cm,
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "Title",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=26,
        leading=30,
        textColor=colors.HexColor(BRAND_DARK),
        spaceAfter=8,
    )
    body_style = ParagraphStyle(
        "Body",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=13,
        leading=18,
        textColor=colors.HexColor(BRAND_DARK),
        spaceAfter=8,
    )
    src_style = ParagraphStyle(
        "Src",
        parent=styles["BodyText"],
        fontName="Helvetica-Oblique",
        fontSize=9,
        leading=11,
        textColor=colors.HexColor("#5b6574"),
    )

    story = []
    for idx, slide in enumerate(SLIDES):
        story.append(Paragraph(slide.title, title_style))
        story.append(Paragraph(slide.paragraph, body_style))
        story.append(Spacer(1, 0.2 * cm))
        if slide.chart_key:
            story.append(Image(str(chart_paths[slide.chart_key]), width=20.4 * cm, height=8.8 * cm))
            story.append(Spacer(1, 0.2 * cm))
        story.append(Paragraph(slide.sources, src_style))
        if idx != len(SLIDES) - 1:
            story.append(PageBreak())
    doc.build(story)
    return out_pdf


def _add_textbox(slide, left: float, top: float, width: float, height: float, text: str, size: int, bold: bool = False) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold


def build_pptx(chart_paths: dict[str, Path]) -> Path | None:
    if not _pptx_available():
        print("[WARN] python-pptx not installed; skipping .pptx generation.")
        print("       Install with: python3 -m pip install python-pptx")
        return None

    OUT_PPTX_DIR.mkdir(parents=True, exist_ok=True)
    out_pptx = OUT_PPTX_DIR / f"investor-presentation-{date.today().isoformat()}.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for s in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _add_textbox(slide, 0.55, 0.30, 12.2, 0.95, s.title, 34, bold=True)
        _add_textbox(slide, 0.65, 1.20, 12.0, 1.55, s.paragraph, 17, bold=False)
        if s.chart_key:
            slide.shapes.add_picture(str(chart_paths[s.chart_key]), Inches(1.0), Inches(2.75), width=Inches(11.4), height=Inches(3.95))
        _add_textbox(slide, 0.65, 6.85, 12.0, 0.45, s.sources, 10, bold=False)

    prs.save(str(out_pptx))
    return out_pptx


def main() -> None:
    chart_paths = _build_chart_files()
    out_pdf = build_pdf(chart_paths)
    out_pptx = build_pptx(chart_paths)
    OUT_NOTES_DIR.mkdir(parents=True, exist_ok=True)
    notes_path = OUT_NOTES_DIR / f"investor-presentation-notes-{date.today().isoformat()}.md"
    notes_lines = ["# Investor Presentation Notes", ""]
    for idx, slide in enumerate(SLIDES, start=1):
        notes_lines.append(f"## {idx}. {slide.title}")
        notes_lines.append(slide.paragraph)
        notes_lines.append("")
        notes_lines.append(f"_Evidence:_ {slide.sources}")
        notes_lines.append("")
    notes_path.write_text("\n".join(notes_lines), encoding="utf-8")
    print(f"[OK] PDF generated: {out_pdf}")
    if out_pptx:
        print(f"[OK] PPTX generated: {out_pptx}")
        print("[INFO] Upload the PPTX to Google Drive and open with Google Slides.")
    print(f"[OK] Notes generated: {notes_path}")


if __name__ == "__main__":
    main()
